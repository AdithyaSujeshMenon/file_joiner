import os
import logging
from PyPDF2 import PdfMerger
from docx import Document
from docxcompose.composer import Composer
import pandas as pd
from markdown2 import markdown
from pptx import Presentation
from zipfile import ZipFile
from shutil import copyfile
from tqdm import tqdm
from odf.opendocument import OpenDocumentText, OpenDocumentSpreadsheet, OpenDocumentPresentation
from odf.text import P
from odf.table import Table, TableRow, TableCell
from odf.draw import Page

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def join_pdfs(files, output_path):
    try:
        merger = PdfMerger()
        for pdf in tqdm(files, desc="Joining PDFs"):
            merger.append(pdf)
        merger.write(output_path)
        merger.close()
        logging.info(f"PDF files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining PDF files: {e}")

def join_docx(files, output_path):
    try:
        master = Document(files[0])
        composer = Composer(master)
        for doc in tqdm(files[1:], desc="Joining DOCX files"):
            doc = Document(doc)
            composer.append(doc)
        composer.save(output_path)
        logging.info(f"DOCX files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining DOCX files: {e}")

def join_txt(files, output_path):
    try:
        with open(output_path, 'w') as outfile:
            for fname in tqdm(files, desc="Joining TXT files"):
                with open(fname) as infile:
                    outfile.write(infile.read())
                    outfile.write("\n")
        logging.info(f"TXT files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining TXT files: {e}")

def join_csv(files, output_path):
    try:
        combined_csv = pd.concat([pd.read_csv(f) for f in tqdm(files, desc="Joining CSV files")])
        combined_csv.to_csv(output_path, index=False)
        logging.info(f"CSV files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining CSV files: {e}")

def join_xlsx(files, output_path):
    try:
        combined_xlsx = pd.concat([pd.read_excel(f) for f in tqdm(files, desc="Joining XLSX files")])
        combined_xlsx.to_excel(output_path, index=False)
        logging.info(f"XLSX files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining XLSX files: {e}")

def join_pptx(files, output_path):
    try:
        presentation = Presentation()
        for file in tqdm(files, desc="Joining PPTX files"):
            temp_pptx = Presentation(file)
            for slide in temp_pptx.slides:
                slide_layout = presentation.slide_layouts[0]
                slide_copy = presentation.slides.add_slide(slide_layout)
                for shape in slide.shapes:
                    new_shape = slide_copy.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
                    new_shape.text = shape.text
        presentation.save(output_path)
        logging.info(f"PPTX files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining PPTX files: {e}")

def join_rtf(files, output_path):
    try:
        with open(output_path, 'w') as outfile:
            for file in tqdm(files, desc="Joining RTF files"):
                with open(file) as infile:
                    outfile.write(infile.read())
                    outfile.write("\n")
        logging.info(f"RTF files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining RTF files: {e}")

def join_epub(files, output_path):
    try:
        with ZipFile(output_path, 'w') as zipf:
            for file in tqdm(files, desc="Joining EPUB files"):
                zipf.write(file, os.path.basename(file))
        logging.info(f"EPUB files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining EPUB files: {e}")

def join_odt(files, output_path):
    try:
        final_doc = OpenDocumentText()
        for file in tqdm(files, desc="Joining ODT files"):
            doc = OpenDocumentText(file)
            for element in doc.body.childNodes:
                final_doc.body.appendChild(element)
        final_doc.save(output_path)
        logging.info(f"ODT files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining ODT files: {e}")

def join_ods(files, output_path):
    try:
        final_doc = OpenDocumentSpreadsheet()
        for file in tqdm(files, desc="Joining ODS files"):
            doc = OpenDocumentSpreadsheet(file)
            for element in doc.body.childNodes:
                final_doc.body.appendChild(element)
        final_doc.save(output_path)
        logging.info(f"ODS files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining ODS files: {e}")

def join_odp(files, output_path):
    try:
        final_doc = OpenDocumentPresentation()
        for file in tqdm(files, desc="Joining ODP files"):
            doc = OpenDocumentPresentation(file)
            for element in doc.body.childNodes:
                final_doc.body.appendChild(element)
        final_doc.save(output_path)
        logging.info(f"ODP files successfully joined into {output_path}")
    except Exception as e:
        logging.error(f"Error joining ODP files: {e}")

def handle_zip(files, output_path):
    try:
        with ZipFile(output_path, 'w') as zipf:
            for file in tqdm(files, desc="Compressing files into ZIP"):
                zipf.write(file, os.path.basename(file))
        logging.info(f"Files successfully compressed into {output_path}")
    except Exception as e:
        logging.error(f"Error compressing files: {e}")

def main():
    file_types = ['pdf', 'docx', 'txt', 'csv', 'xlsx', 'pptx', 'rtf', 'epub', 'odt', 'ods', 'odp', 'zip']
    
    print("Welcome to the enhanced file joiner program!")
    
    file_type = input(f"Enter the file type to join ({', '.join(file_types)}): ").strip().lower()
    if file_type not in file_types:
        print(f"Unsupported file type: {file_type}")
        return

    try:
        num_files = int(input("Enter the number of files to join: "))
        if num_files <= 0:
            raise ValueError("The number of files must be greater than zero.")
    except ValueError as e:
        print(f"Invalid number of files: {e}")
        return
    
    files = []
    for i in range(num_files):
        file_path = input(f"Enter the path for file {i+1}: ").strip()
        if not os.path.isfile(file_path):
            print(f"File not found: {file_path}")
            return
        files.append(file_path)
    
    # Display the list of files and their current order
    print("Here are the files you have selected:")
    for idx, file in enumerate(files, 1):
        print(f"{idx}. {file}")
    
    # Ask for the desired order
    print("Please specify the new order of the files by entering the numbers separated by commas (e.g., 2,1 for reversing the order).")
    try:
        new_order = list(map(int, input("Enter the new order: ").strip().split(',')))
        if len(new_order) != num_files or any(o < 1 or o > num_files for o in new_order):
            raise ValueError("Invalid order specified.")
        ordered_files = [files[i-1] for i in new_order]
    except ValueError as e:
        print(f"Invalid order input: {e}")
        return

    output_file = input(f"Enter the output file name (including .{file_type} extension): ").strip()

    print("\nSummary of Actions:")
    print(f"File type: {file_type}")
    print(f"Number of files: {num_files}")
    print("Order of files:")
    for idx, file in enumerate(ordered_files, 1):
        print(f"{idx}. {file}")
    print(f"Output file: {output_file}")
    
    confirmation = input("Do you want to proceed with the above actions? (yes/no): ").strip().lower()
    if confirmation != 'yes':
        print("Operation cancelled.")
        return

    try:
        if file_type == 'pdf':
            join_pdfs(ordered_files, output_file)
        elif file_type == 'docx':
            join_docx(ordered_files, output_file)
        elif file_type == 'txt':
            join_txt(ordered_files, output_file)
        elif file_type == 'csv':
            join_csv(ordered_files, output_file)
        elif file_type == 'xlsx':
            join_xlsx(ordered_files, output_file)
        elif file_type == 'pptx':
            join_pptx(ordered_files, output_file)
        
        elif file_type == 'rtf':
            join_rtf(ordered_files, output_file)
        elif file_type == 'epub':
            join_epub(ordered_files, output_file)
        elif file_type == 'odt':
            join_odt(ordered_files, output_file)
        elif file_type == 'ods':
            join_ods(ordered_files, output_file)
        elif file_type == 'odp':
            join_odp(ordered_files, output_file)
        elif file_type == 'zip':
            handle_zip(ordered_files, output_file)
        else:
            print(f"Unsupported file type: {file_type}")
    except Exception as e:
        print(f"An error occurred during file joining: {e}")

if __name__ == "__main__":
    main()
